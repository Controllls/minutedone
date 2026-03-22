"""
app.py — 티로 회의록 웹 인터페이스
실행: python app.py
접속: http://localhost:5000
"""

import os
import sys
import json
import glob as glob_module
from flask import Flask, render_template, request, jsonify, redirect
from dotenv import load_dotenv

# localhost HTTP에서 OAuth 허용
os.environ.setdefault("OAUTHLIB_INSECURE_TRANSPORT", "1")

load_dotenv("config/.env")
load_dotenv("config/calendar.env", override=False)

sys.path.insert(0, "scripts")
import stt_calendar
import rag
import notion as notion_module
import db
import llm as llm_module
import rules as rules_module

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # 100MB
db.init_db()


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/run-calendar", methods=["POST"])
def run_calendar():
    data = request.get_json()
    transcript = data.get("transcript", "").strip()
    if not transcript:
        return jsonify({"error": "녹취록 텍스트가 없습니다."}), 400

    company     = data.get("company", "").strip() or None
    ppt_context = data.get("ppt_context", "").strip() or None
    try:
        events = stt_calendar.extract_calendar_events(transcript, company=company, ppt_context=ppt_context)
        created = stt_calendar.create_notion_calendar_events(events)
        return jsonify({"success": True, "events": events, "created": len(created), "usage": llm_module.get_usage()})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/extract", methods=["POST"])
def extract_only():
    """캘린더 등록 없이 추출만 (미리보기용)"""
    data = request.get_json()
    transcript = data.get("transcript", "").strip()
    if not transcript:
        return jsonify({"error": "녹취록 텍스트가 없습니다."}), 400

    company     = data.get("company", "").strip() or None
    ppt_context = data.get("ppt_context", "").strip() or None
    try:
        events   = stt_calendar.extract_calendar_events(transcript, company=company, ppt_context=ppt_context)
        insights = stt_calendar.extract_insights(transcript)
        return jsonify({"success": True, "events": events, "insights": insights, "usage": llm_module.get_usage()})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/chat", methods=["POST"])
def chat():
    data = request.get_json()
    question = data.get("question", "").strip()
    history  = data.get("history", [])
    if not question:
        return jsonify({"error": "질문이 없습니다."}), 400
    try:
        reply, updated_history = rag.chat(question, history)
        return jsonify({"success": True, "reply": reply, "history": updated_history})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/rag/add", methods=["POST"])
def rag_add():
    """회의록을 RAG 저장소에 추가"""
    data = request.get_json()
    text  = data.get("text", "").strip()
    title = data.get("title", "회의록")
    if not text:
        return jsonify({"error": "텍스트가 없습니다."}), 400
    chunks = rag.add_document(text, metadata={"title": title})
    return jsonify({"success": True, "chunks_added": chunks, "stats": rag.store_stats()})


@app.route("/api/rag/stats", methods=["GET"])
def rag_stats():
    return jsonify(rag.store_stats())


@app.route("/api/rag/documents", methods=["GET"])
def rag_documents():
    docs = rag.get_all_documents()
    return jsonify({"success": True, "documents": docs, "total": len(docs)})


@app.route("/api/rag/sync-notion", methods=["POST"])
def rag_sync_notion():
    """노션 회의록 데이터베이스를 RAG 저장소에 동기화"""
    data = request.get_json() or {}
    database_id = data.get("database_id") or os.getenv("NOTION_MEETINGS_DB_ID", "").strip()
    clear_first = data.get("clear_first", False)

    if not database_id:
        return jsonify({"error": "database_id 가 없습니다. 요청 본문 또는 NOTION_MEETINGS_DB_ID 환경변수를 확인하세요."}), 400

    try:
        meetings = notion_module.get_meetings_for_rag(database_id)
        if not meetings:
            return jsonify({"success": True, "synced": 0, "chunks_added": 0, "stats": rag.store_stats()})

        if clear_first:
            rag.clear_store()

        total_chunks = 0
        for m in meetings:
            added = rag.add_document(
                m["text"],
                metadata={"title": m["title"], "notion_page_id": m["page_id"], "date": m["created_time"]},
            )
            total_chunks += added

        return jsonify({
            "success": True,
            "synced": len(meetings),
            "chunks_added": total_chunks,
            "stats": rag.store_stats(),
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/state/save", methods=["POST"])
def state_save():
    data = request.get_json() or {}
    events     = data.get("events", [])
    task_views = data.get("taskViews", {})
    recent     = data.get("recentList", [])

    db.save_events(events)
    db.kv_set("taskViews", task_views)
    for m in recent:
        pass  # recentList는 add_recent_meeting으로 개별 저장
    return jsonify({"success": True})


@app.route("/api/state/load", methods=["GET"])
def state_load():
    return jsonify({
        "success":    True,
        "events":     db.load_events(),
        "taskViews":  db.kv_get("taskViews", {}),
        "recentList": db.load_recent_meetings(),
    })


@app.route("/api/state/meeting", methods=["POST"])
def state_add_meeting():
    """최근 회의 1건 저장"""
    data = request.get_json() or {}
    db.add_recent_meeting(
        preview    = data.get("preview", ""),
        time       = data.get("time", ""),
        task_count = data.get("count", 0),
    )
    return jsonify({"success": True})


@app.route("/api/classify-tasks", methods=["POST"])
def classify_tasks():
    data = request.get_json()
    events = data.get("events", [])
    user_name = data.get("user_name", "전체")
    if not events:
        return jsonify({"error": "events 가 없습니다."}), 400
    try:
        result = stt_calendar.classify_tasks(events, user_name)
        return jsonify({"success": True, "views": result.get("views", {})})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/rag/clear", methods=["POST"])
def rag_clear():
    """RAG 저장소와 DB 이벤트를 모두 초기화합니다."""
    rag.clear_store()
    db.save_events([])
    db.kv_set("taskViews", {})
    return jsonify({"success": True})


@app.route("/api/stt", methods=["POST"])
def stt():
    """Whisper STT: 업로드된 오디오 파일을 텍스트로 변환합니다."""
    if "audio" not in request.files:
        return jsonify({"error": "audio 파일이 없습니다."}), 400

    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    if not api_key:
        return jsonify({"error": "OPENAI_API_KEY 가 설정되지 않았습니다. 브라우저 STT 모드를 사용하세요."}), 400

    audio_file = request.files["audio"]
    try:
        from openai import OpenAI
        import traceback
        client = OpenAI(api_key=api_key)
        audio_bytes = audio_file.read()
        filename = audio_file.filename or "audio.webm"
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "webm"
        print(f"[STT] filename={filename!r} ext={ext} size={len(audio_bytes)} mime={audio_file.mimetype} magic={audio_bytes[:12].hex()}")

        import io

        # 3gp4 컨테이너(삼성 통화 녹음 등)는 ftyp 브랜드를 isom으로 패치 후 .mp4로 전달
        ftyp_brand = audio_bytes[8:12].decode("latin-1", errors="replace") if len(audio_bytes) >= 12 else ""
        print(f"[STT] ftyp_brand={ftyp_brand!r}")

        if ftyp_brand.startswith("3gp") or ext in ("3gp", "3gpp"):
            patched = bytearray(audio_bytes)
            patched[8:12] = b"isom"  # major brand 패치
            buf = io.BytesIO(bytes(patched))
            buf.name = filename.rsplit(".", 1)[0] + ".mp4"
            print(f"[STT] 3gp4→mp4 패치 적용, new name={buf.name}")
        else:
            buf = io.BytesIO(audio_bytes)
            buf.name = filename

        _WHISPER_PROMPT = (
            "DF 기획팀: 김동석 프로, 김현수 프로, 남예서 프로, 옥영빈 프로, 김희주 프로, 이하은 팀장, 유형욱 프로, 이현주 이사님. "
            "DF 디자인팀: 김득환 프로, 강수빈 프로, 김민경 프로, 이주희 프로, 이주환 이사님. "
            "DF 모션팀: 유근주 실장님, 김유미 프로, 최유정 프로. "
            "IKEA: 노미소님, 세레나님, 수아님, 제니님, 인국님, 티프님, 신디님, 루카님, 크리스틴님. "
            "협력사: 보배 카피님, 은선 카피님, 다영 차장님, 미소 대리님."
        )
        transcript = client.audio.transcriptions.create(
            model="whisper-1",
            file=buf,
            language="ko",
            prompt=_WHISPER_PROMPT,
        )
        return jsonify({"success": True, "text": transcript.text})
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/usage", methods=["GET"])
def usage():
    return jsonify(llm_module.get_usage())


@app.route("/api/rules", methods=["GET"])
def get_rules():
    """사용 가능한 규칙 파일 목록과 현재 규칙 내용을 반환합니다."""
    company = request.args.get("company", "").strip() or None
    return jsonify({
        "available": rules_module.list_available(),
        "current_company": company or os.getenv("COMPANY_RULES", ""),
        "content": rules_module.load(company),
    })


@app.route("/api/rules/<company>", methods=["GET"])
def get_rule_file(company):
    """특정 회사 규칙 파일 내용을 반환합니다."""
    from pathlib import Path
    path = Path(f"config/rules/{company}.md")
    if not path.exists():
        return jsonify({"error": f"{company}.md 파일이 없습니다."}), 404
    return jsonify({"company": company, "content": path.read_text(encoding="utf-8")})


@app.route("/api/rules/<company>", methods=["POST"])
def save_rule_file(company):
    """회사 규칙 파일을 저장합니다."""
    from pathlib import Path
    data = request.get_json() or {}
    content = data.get("content", "")
    path = Path(f"config/rules/{company}.md")
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return jsonify({"success": True, "company": company})


_PPT_DIR = os.path.join(os.path.dirname(__file__), "uploads", "ppt")
os.makedirs(_PPT_DIR, exist_ok=True)

_PPT_EXTS = {".ppt", ".pptx", ".pdf", ".key", ".odp"}


@app.route("/api/ppt/list", methods=["GET"])
def ppt_list():
    files = []
    for f in sorted(os.listdir(_PPT_DIR), key=lambda x: os.path.getmtime(os.path.join(_PPT_DIR, x)), reverse=True):
        path = os.path.join(_PPT_DIR, f)
        if os.path.isfile(path):
            stat = os.stat(path)
            files.append({
                "name": f,
                "size": stat.st_size,
                "mtime": stat.st_mtime,
            })
    return jsonify({"success": True, "files": files})


@app.route("/api/ppt/upload", methods=["POST"])
def ppt_upload():
    if "file" not in request.files:
        return jsonify({"error": "파일이 없습니다."}), 400
    file = request.files["file"]
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in _PPT_EXTS:
        return jsonify({"error": f"지원하지 않는 형식입니다. ({', '.join(_PPT_EXTS)})"}), 400
    save_path = os.path.join(_PPT_DIR, file.filename)
    file.save(save_path)
    return jsonify({"success": True, "name": file.filename})


@app.route("/api/ppt/extract/<filename>", methods=["GET"])
def ppt_extract(filename):
    """PPT/PDF 파일에서 텍스트를 추출합니다."""
    path = os.path.join(_PPT_DIR, filename)
    if not os.path.isfile(path):
        return jsonify({"error": "파일을 찾을 수 없습니다."}), 404
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext in (".ppt", ".pptx"):
            from pptx import Presentation
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                slide_lines = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_lines.append(text)
                if slide_lines:
                    lines.append(f"[슬라이드 {i}]\n" + "\n".join(slide_lines))
            text = "\n\n".join(lines)
        elif ext == ".pdf":
            import pdfplumber
            pages = []
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    t = (page.extract_text() or "").strip()
                    if t:
                        pages.append(f"[페이지 {i}]\n{t}")
            text = "\n\n".join(pages)
        else:
            return jsonify({"error": "텍스트 추출을 지원하지 않는 형식입니다."}), 400

        return jsonify({"success": True, "filename": filename, "text": text, "length": len(text)})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/ppt/download/<filename>", methods=["GET"])
def ppt_download(filename):
    from flask import send_from_directory
    return send_from_directory(_PPT_DIR, filename, as_attachment=True)


@app.route("/api/ppt/delete/<filename>", methods=["DELETE"])
def ppt_delete(filename):
    path = os.path.join(_PPT_DIR, filename)
    if not os.path.isfile(path):
        return jsonify({"error": "파일을 찾을 수 없습니다."}), 404
    os.remove(path)
    return jsonify({"success": True})


# ── Google Drive ──────────────────────────────────────────────
_GDRIVE_TOKEN   = os.path.join(os.path.dirname(__file__), "config", "gdrive_token.json")
_GDRIVE_SCOPES  = ["https://www.googleapis.com/auth/drive.readonly"]
_GDRIVE_REDIRECT = "http://localhost:5000/api/gdrive/callback"


def _gdrive_creds():
    """저장된 토큰 로드. 만료 시 자동 갱신."""
    if not os.path.exists(_GDRIVE_TOKEN):
        return None
    try:
        from google.oauth2.credentials import Credentials
        from google.auth.transport.requests import Request as GRequest
        import traceback
        creds = Credentials.from_authorized_user_file(_GDRIVE_TOKEN, _GDRIVE_SCOPES)
        if not creds.valid:
            if creds.refresh_token:
                try:
                    creds.refresh(GRequest())
                    open(_GDRIVE_TOKEN, "w").write(creds.to_json())
                except Exception as ref_err:
                    print(f"[gdrive] 토큰 갱신 실패: {ref_err}")
                    traceback.print_exc()
                    return None
            else:
                print("[gdrive] refresh_token 없음 — 재인증 필요")
                return None
        return creds if creds.valid else None
    except Exception:
        import traceback
        traceback.print_exc()
        return None


def _gdrive_client_config():
    client_id     = os.getenv("GOOGLE_CLIENT_ID", "").strip()
    client_secret = os.getenv("GOOGLE_CLIENT_SECRET", "").strip()
    if not client_id or not client_secret:
        return None
    return {"web": {
        "client_id":     client_id,
        "client_secret": client_secret,
        "auth_uri":      "https://accounts.google.com/o/oauth2/auth",
        "token_uri":     "https://oauth2.googleapis.com/token",
    }}


def _gdrive_flow(state=None):
    from google_auth_oauthlib.flow import Flow
    cfg = _gdrive_client_config()
    if not cfg:
        return None
    return Flow.from_client_config(
        cfg,
        scopes=_GDRIVE_SCOPES,
        redirect_uri=_GDRIVE_REDIRECT,
        state=state,
    )


@app.route("/api/gdrive/status")
def gdrive_status():
    creds = _gdrive_creds()
    return jsonify({"connected": creds is not None})


@app.route("/api/gdrive/auth")
def gdrive_auth():
    cfg = _gdrive_client_config()
    if not cfg:
        return jsonify({"error": "GOOGLE_CLIENT_ID / GOOGLE_CLIENT_SECRET 를 config/.env 에 설정하세요."}), 400
    import urllib.parse, secrets
    state = secrets.token_urlsafe(24)
    open("config/gdrive_state.tmp", "w").write(state)
    params = {
        "client_id":     cfg["web"]["client_id"],
        "redirect_uri":  _GDRIVE_REDIRECT,
        "response_type": "code",
        "scope":         " ".join(_GDRIVE_SCOPES),
        "access_type":   "offline",
        "prompt":        "consent",
        "state":         state,
    }
    return redirect("https://accounts.google.com/o/oauth2/auth?" + urllib.parse.urlencode(params))


@app.route("/api/gdrive/callback")
def gdrive_callback():
    cfg  = _gdrive_client_config()
    if not cfg:
        return "설정 오류", 400
    code = request.args.get("code")
    if not code:
        return f"code 없음: {dict(request.args)}", 400
    try:
        import requests as _req
        resp = _req.post("https://oauth2.googleapis.com/token", data={
            "code":          code,
            "client_id":     cfg["web"]["client_id"],
            "client_secret": cfg["web"]["client_secret"],
            "redirect_uri":  _GDRIVE_REDIRECT,
            "grant_type":    "authorization_code",
        })
        if not resp.ok:
            return f"토큰 교환 실패: {resp.text}", 400
        tok = resp.json()
        from google.oauth2.credentials import Credentials
        creds = Credentials(
            token=tok.get("access_token"),
            refresh_token=tok.get("refresh_token"),
            token_uri="https://oauth2.googleapis.com/token",
            client_id=cfg["web"]["client_id"],
            client_secret=cfg["web"]["client_secret"],
            scopes=_GDRIVE_SCOPES,
        )
        open(_GDRIVE_TOKEN, "w").write(creds.to_json())
        return redirect("/?gdrive=connected")
    except Exception as e:
        return f"인증 실패: {e}", 400


@app.route("/api/gdrive/browse")
def gdrive_browse():
    """폴더 내용 탐색 — folder_id 없으면 root"""
    creds = _gdrive_creds()
    if not creds:
        return jsonify({"error": "인증이 필요합니다."}), 401
    folder_id = request.args.get("folder_id", "root")
    try:
        from googleapiclient.discovery import build
        svc = build("drive", "v3", credentials=creds)

        # 현재 폴더 이름
        if folder_id == "root":
            folder_name = "내 드라이브"
            parent_id   = None
        else:
            meta = svc.files().get(fileId=folder_id, fields="id,name,parents").execute()
            folder_name = meta.get("name", folder_id)
            parent_id   = (meta.get("parents") or [None])[0]

        # 하위 항목
        res = svc.files().list(
            q=f"'{folder_id}' in parents and trashed = false",
            pageSize=100,
            fields="files(id,name,mimeType,size,modifiedTime)",
            orderBy="folder,name",
        ).execute()

        items = res.get("files", [])
        folders = [f for f in items if f["mimeType"] == "application/vnd.google-apps.folder"]
        files   = [f for f in items if f["name"].lower().endswith(".mp3")]

        return jsonify({
            "success":     True,
            "folder_id":   folder_id,
            "folder_name": folder_name,
            "parent_id":   parent_id,
            "folders":     folders,
            "files":       files,
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/gdrive/files")
def gdrive_files():
    creds = _gdrive_creds()
    if not creds:
        return jsonify({"error": "인증이 필요합니다."}), 401
    folder_id = request.args.get("folder_id", "root")
    try:
        from googleapiclient.discovery import build
        svc   = build("drive", "v3", credentials=creds)
        parent = f"'{folder_id}' in parents" if folder_id else "true"
        # 오디오 파일 + 흔한 녹음 확장자 모두 포함
        query = (
            f"({parent}) and trashed = false and ("
            "mimeType contains 'audio/' or "
            "name contains '.m4a' or name contains '.mp3' or "
            "name contains '.wav' or name contains '.ogg' or "
            "name contains '.flac' or name contains '.aac' or "
            "name contains '.amr' or name contains '.3gp'"
            ")"
        )
        res = svc.files().list(
            q=query, pageSize=100,
            fields="files(id,name,size,modifiedTime,mimeType)",
            orderBy="modifiedTime desc",
        ).execute()
        return jsonify({"success": True, "files": res.get("files", [])})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/gdrive/transcribe/<file_id>", methods=["POST"])
def gdrive_transcribe(file_id):
    creds = _gdrive_creds()
    if not creds:
        return jsonify({"error": "인증이 필요합니다."}), 401
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    if not api_key:
        return jsonify({"error": "OPENAI_API_KEY 가 설정되지 않았습니다."}), 400
    try:
        import io
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaIoBaseDownload
        svc  = build("drive", "v3", credentials=creds)
        meta = svc.files().get(fileId=file_id, fields="name,size,mimeType").execute()
        name = meta["name"]
        mime = meta.get("mimeType", "audio/mpeg")
        size = int(meta.get("size", 0))
        if size > 25 * 1024 * 1024:
            return jsonify({"error": f"파일이 25MB를 초과합니다. ({size//1024//1024}MB)"}), 400

        buf = io.BytesIO()
        dl  = MediaIoBaseDownload(buf, svc.files().get_media(fileId=file_id))
        done = False
        while not done:
            _, done = dl.next_chunk()
        buf.seek(0)

        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        _WHISPER_PROMPT = (
            "DF 기획팀: 김동석 프로, 김현수 프로, 남예서 프로, 옥영빈 프로, 김희주 프로, 이하은 팀장, 유형욱 프로, 이현주 이사님. "
            "DF 디자인팀: 김득환 프로, 강수빈 프로, 김민경 프로, 이주희 프로, 이주환 이사님. "
            "DF 모션팀: 유근주 실장님, 김유미 프로, 최유정 프로. "
            "IKEA: 노미소님, 세레나님, 수아님, 제니님, 인국님, 티프님, 신디님, 루카님, 크리스틴님. "
            "협력사: 보배 카피님, 은선 카피님, 다영 차장님, 미소 대리님."
        )
        transcript = client.audio.transcriptions.create(
            model="whisper-1", file=(name, buf, mime), language="ko", prompt=_WHISPER_PROMPT
        )
        return jsonify({"success": True, "text": transcript.text, "filename": name})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/gdrive/disconnect", methods=["POST"])
def gdrive_disconnect():
    if os.path.exists(_GDRIVE_TOKEN):
        os.remove(_GDRIVE_TOKEN)
    return jsonify({"success": True})


# ============================================================
# 관리자 — Rules
# ============================================================

_RULES_DIR = os.path.join(os.path.dirname(__file__), "config", "rules")


def _safe_rule_name(name):
    """파일 이름 경로 traversal 방지"""
    return "/" not in name and ".." not in name and name.endswith(".md")


@app.route("/api/admin/rules", methods=["GET"])
def admin_list_rules():
    os.makedirs(_RULES_DIR, exist_ok=True)
    files = sorted(glob_module.glob(os.path.join(_RULES_DIR, "*.md")))
    return jsonify([os.path.basename(f) for f in files])


@app.route("/api/admin/rules/<name>", methods=["GET"])
def admin_get_rule(name):
    if not _safe_rule_name(name):
        return jsonify({"error": "invalid"}), 400
    path = os.path.join(_RULES_DIR, name)
    if not os.path.exists(path):
        return jsonify({"error": "not found"}), 404
    with open(path, "r", encoding="utf-8") as f:
        return jsonify({"name": name, "content": f.read()})


@app.route("/api/admin/rules/<name>", methods=["POST"])
def admin_save_rule(name):
    if not _safe_rule_name(name):
        return jsonify({"error": "invalid name — .md 파일만 허용"}), 400
    content = (request.get_json() or {}).get("content", "")
    os.makedirs(_RULES_DIR, exist_ok=True)
    with open(os.path.join(_RULES_DIR, name), "w", encoding="utf-8") as f:
        f.write(content)
    return jsonify({"ok": True})


@app.route("/api/admin/rules/<name>", methods=["DELETE"])
def admin_delete_rule(name):
    if not _safe_rule_name(name):
        return jsonify({"error": "invalid"}), 400
    path = os.path.join(_RULES_DIR, name)
    if os.path.exists(path):
        os.remove(path)
    return jsonify({"ok": True})


# ============================================================
# 관리자 — PPT 파일 목록 (업로드/삭제는 기존 /api/ppt/* 재사용)
# ============================================================

if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5000)
